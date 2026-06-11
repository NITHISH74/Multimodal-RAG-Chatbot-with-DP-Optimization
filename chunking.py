"""
Document parsing + chunking (Phase 2).

Turns an uploaded file into a list of clean, deduplicated chunk dicts ready
for embedding. Splitting respects page/slide boundaries and prefers
paragraph- then sentence-aware breaks. Empty / garbage chunks are filtered
out and each surviving chunk carries a content hash for DB-level dedup.

A chunk dict looks like::

    {
        "content": "...",          # text to embed (or "" for images)
        "file_name": "report.pdf",
        "document_type": "pdf",     # pdf|docx|pptx|text|image|web
        "page_number": 3,           # page (pdf) / slide (pptx) / None
        "chunk_index": 0,           # ordinal within the document
        "content_hash": "ab12...",  # sha256(file_name + content)
        "type": "text",             # text|image  (drives embedding path)
    }
"""
import io
import re
import hashlib

from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

import config

_SENTENCE_RE = re.compile(r"(?<=[.!?。！？])\s+")
_WORD_RE = re.compile(r"[A-Za-z0-9À-￿]")


# ──────────────────────────────────────────────────────────────────────
#  Cleaning / quality filtering
# ──────────────────────────────────────────────────────────────────────
def normalize_text(text):
    """Collapse runaway whitespace while preserving paragraph breaks."""
    if not text:
        return ""
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def is_meaningful(text, min_chars=None):
    """Reject empty, too-short, or gibberish chunks (Phase 2.4).

    A chunk is meaningful when, after stripping, it has at least
    ``min_chars`` characters AND a reasonable fraction of them are
    word-characters (filters out lines of punctuation / control noise).
    """
    if min_chars is None:
        min_chars = config.MIN_CHUNK_CHARS
    if not text:
        return False
    stripped = text.strip()
    if len(stripped) < min_chars:
        return False
    word_chars = len(_WORD_RE.findall(stripped))
    # At least 40% of the content should be actual word characters.
    return word_chars >= max(min_chars * 0.4, 0.4 * len(stripped) * 0.5)


def content_hash(file_name, content):
    """Stable hash of (file_name + normalized content) for dedup."""
    key = f"{file_name}::{normalize_text(content)}".encode("utf-8")
    return hashlib.sha256(key).hexdigest()


# ──────────────────────────────────────────────────────────────────────
#  Sentence / paragraph aware splitting
# ──────────────────────────────────────────────────────────────────────
def split_text(text, target_chars=None, overlap_chars=None):
    """Split text into chunks, preferring paragraph then sentence breaks.

    Greedily packs paragraphs up to ``target_chars``; oversized paragraphs
    are further split on sentence boundaries. A small character overlap is
    carried between consecutive chunks to preserve context across breaks.
    """
    target = target_chars or config.CHUNK_TARGET_CHARS
    overlap = overlap_chars if overlap_chars is not None else config.CHUNK_OVERLAP_CHARS
    text = normalize_text(text)
    if not text:
        return []

    # Break into sentence-level units first (via paragraphs).
    units = []
    for para in re.split(r"\n{2,}", text):
        para = para.strip()
        if not para:
            continue
        if len(para) <= target:
            units.append(para)
        else:
            units.extend(s.strip() for s in _SENTENCE_RE.split(para) if s.strip())

    chunks, buf = [], ""
    for unit in units:
        if len(unit) > target:
            # A single very long unit: hard-wrap it.
            if buf:
                chunks.append(buf)
                buf = ""
            for i in range(0, len(unit), target):
                chunks.append(unit[i:i + target])
            continue
        if buf and len(buf) + len(unit) + 1 > target:
            chunks.append(buf)
            tail = buf[-overlap:] if overlap else ""
            buf = (tail + " " + unit).strip()
        else:
            buf = (buf + " " + unit).strip() if buf else unit
    if buf:
        chunks.append(buf)
    return chunks


# ──────────────────────────────────────────────────────────────────────
#  Per-format extraction → page/slide-tagged text segments
# ──────────────────────────────────────────────────────────────────────
def _segments_from_pdf(file_bytes):
    reader = PdfReader(io.BytesIO(file_bytes))
    for page_no, page in enumerate(reader.pages, start=1):
        yield page_no, page.extract_text() or ""


def _segments_from_docx(file_bytes):
    # DOCX has no reliable page concept; treat the whole document as page 1
    # but keep paragraph structure (blank lines) for paragraph-aware splitting.
    doc = DocxDocument(io.BytesIO(file_bytes))
    text = "\n\n".join(p.text for p in doc.paragraphs if p.text and p.text.strip())
    yield None, text


def _segments_from_pptx(file_bytes):
    prs = Presentation(io.BytesIO(file_bytes))
    for slide_no, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
        yield slide_no, "\n".join(parts)


def _segments_from_text(file_bytes):
    yield None, file_bytes.decode("utf-8", errors="ignore")


_EXTRACTORS = {
    "pdf": _segments_from_pdf,
    "docx": _segments_from_docx,
    "pptx": _segments_from_pptx,
    "txt": _segments_from_text,
    "md": _segments_from_text,
}


def detect_type(file_name):
    ext = file_name.rsplit(".", 1)[-1].lower() if "." in file_name else ""
    if ext in config.ALLOWED_IMG_EXTS:
        return ext, "image"
    if ext in ("txt", "md"):
        return ext, "text"
    return ext, ext  # pdf/docx/pptx -> document_type == ext


# ──────────────────────────────────────────────────────────────────────
#  Public entry points
# ──────────────────────────────────────────────────────────────────────
def chunk_document(file_name, file_bytes):
    """Parse + chunk a text-bearing document into clean chunk dicts."""
    ext, dtype = detect_type(file_name)
    extractor = _EXTRACTORS.get(ext)
    if extractor is None:
        return []

    chunks, idx, seen = [], 0, set()
    for page_no, segment_text in extractor(file_bytes):
        for piece in split_text(segment_text):
            if not is_meaningful(piece):
                continue
            h = content_hash(file_name, piece)
            if h in seen:               # in-document dedup
                continue
            seen.add(h)
            chunks.append({
                "content": piece,
                "file_name": file_name,
                "document_type": "text" if dtype == "text" else dtype,
                "page_number": page_no,
                "chunk_index": idx,
                "content_hash": h,
                "type": "text",
            })
            idx += 1
    return chunks


def chunk_web_text(source_url, file_name, cleaned_text, extra_metadata=None):
    """Chunk already-cleaned web text (Phase 10). document_type == 'web'.

    extra_metadata (optional dict) is carried on each chunk and merged into
    the row's metadata jsonb at upsert time (page_title, domain,
    crawl_timestamp, source_type ... for crawled pages).
    """
    chunks, idx, seen = [], 0, set()
    for piece in split_text(cleaned_text):
        if not is_meaningful(piece):
            continue
        h = content_hash(source_url, piece)
        if h in seen:
            continue
        seen.add(h)
        chunk = {
            "content": piece,
            "file_name": file_name,
            "document_type": "web",
            "page_number": None,
            "chunk_index": idx,
            "content_hash": h,
            "type": "text",
            "source_url": source_url,
        }
        if extra_metadata:
            chunk["extra_metadata"] = extra_metadata
        chunks.append(chunk)
        idx += 1
    return chunks


def image_chunk(file_name):
    """A single non-text chunk record for an image."""
    return {
        "content": "",
        "file_name": file_name,
        "document_type": "image",
        "page_number": None,
        "chunk_index": 0,
        "content_hash": content_hash(file_name, "image::" + file_name),
        "type": "image",
    }
